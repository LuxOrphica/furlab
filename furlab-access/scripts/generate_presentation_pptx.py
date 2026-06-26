from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "presentation_assets"
OUT = ROOT / "docs" / "FurLab_DB_Presentation_Draft_2026-03-11.pptx"


def set_run_style(run, size=24, bold=False, color=(20, 30, 55)):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Calibri"


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = subtitle
    set_run_style(r, size=20, color=(70, 80, 110))


def add_content_slide(prs: Presentation, title: str, bullets: list[str], image_name: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_style(r, size=34, bold=True)

    if image_name:
        image_path = ASSETS / image_name
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.0), width=Inches(12.3))
        bullets_top = 6.15
    else:
        bullets_top = 1.25

    body = slide.shapes.add_textbox(Inches(0.8), Inches(bullets_top), Inches(11.8), Inches(1.2 if image_name else 5.6))
    btf = body.text_frame
    btf.clear()
    btf.word_wrap = True
    for idx, line in enumerate(bullets):
        p = btf.paragraphs[0] if idx == 0 else btf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        set_run_style(run, size=20 if image_name else 24, color=(35, 35, 40))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "База данных FurLab",
        f"Черновик презентации • Срез данных на 11 марта 2026 • Сгенерировано {date.today().isoformat()}",
    )

    slides = [
        (
            "1. Зачем нужна БД FurLab",
            [
                "Единый контур учета меховых отходов и сценариев выкладки.",
                "Фиксируем не только текущее состояние, но и историю применения куска в запуске.",
                "Основа для аудита, воспроизводимости и управляемости операций.",
            ],
            None,
        ),
        (
            "2. На чем основана модель",
            [
                "Канон терминов и атрибутов: Приложение X.",
                "Типы выкладки и сценарии выполнения: Приложение Y.",
                "Физическая реализация: Access-схема + SQL-миграции.",
            ],
            None,
        ),
        (
            "3. Архитектура решения",
            [
                "UI и формы работают через API, а не напрямую с БД.",
                "Бизнес-правила и валидации централизованы в сервисном слое.",
            ],
            "01_architecture_overview.png",
        ),
        (
            "4. Каркас предметной модели",
            [
                "Линия геометрии: Part -> Zone -> Fragment -> Layout -> LayoutRun.",
                "Линия склада: ScrapPiece + словари + StorageLocation.",
                "Связка двух линий: LayoutRunScrapPlacement.",
            ],
            "02_er_core.png",
        ),
        (
            "5. Складской контур",
            [
                "ScrapPiece хранит метку, геометрию, качество, статус и метрики куска.",
                "inventoryTag уникален и связывает запись с физическим носителем.",
                "napDirectionDeg фиксирует направление ворса в градусах.",
            ],
            "06_scrappiece_table_sample.png",
        ),
        (
            "6. Операции и статусы",
            [
                "Разрешенные переходы: reserve/release/use по доменным правилам.",
                "Discarded трактуется как терминальное состояние.",
            ],
            "04_status_transitions_rules.png",
        ),
        (
            "7. Типы выкладок по Приложению Y",
            [
                "RegularLayout, IrregularLayout, FillRemainingAreaLayout, InventoryLayout.",
                "Тип фиксируется в Layout.layoutType, параметры — в Layout.params.",
                "Снимок конкретного запуска — LayoutRun.paramsSnapshot.",
            ],
            None,
        ),
        (
            "8. Воспроизводимость результата",
            [
                "Layout.params — рабочая конфигурация.",
                "LayoutRun.paramsSnapshot — зафиксированный состав параметров запуска.",
                "Placement и snapshots позволяют повторить и проверить результат.",
            ],
            "05_traceability_chain.png",
        ),
        (
            "9. Режимы инвентарной выкладки",
            [
                "Режим A: фрагменты как производные от ScrapPiece.",
                "Режим B: назначение ScrapPiece на уже сформированные Fragment.",
                "Оба режима отражаются в LayoutRunScrapPlacement.",
            ],
            None,
        ),
        (
            "10. Фактический срез данных (11.03.2026)",
            [
                "ScrapPiece: 121; статусы: Available 82, Reserved 28, Used 7, Discarded 4.",
                "ScrapTransaction: 34 записи; преобладают Reserve/Release.",
                "LayoutRunScrapPlacement: 7 записей размещения.",
            ],
            "03_status_distribution.png",
        ),
        (
            "11. Эксплуатация и надежность",
            [
                "SchemaMigrations фиксирует примененные изменения схемы (7 миграций).",
                "Есть runbook, backup/restore и smoke/contract тесты API.",
            ],
            "08_migrations_and_table_counts.png",
        ),
        (
            "12. Риски и следующие шаги",
            [
                "Повысить полноту фиксации Use/Discard операций.",
                "Заполнить InventoryLayoutConfig для эксплуатационных сценариев.",
                "Снизить долю ScrapPiece с пустым materialId.",
            ],
            None,
        ),
        (
            "13. Резюме",
            [
                "Модель соответствует Приложению Y и поддерживает трассируемость.",
                "Есть база для аналитики и улучшения алгоритмов выкладки.",
                "Следующий этап — KPI качества данных и операционная дисциплина.",
            ],
            "07_transaction_history_sample.png",
        ),
    ]

    for title, bullets, image in slides:
        add_content_slide(prs, title, bullets, image_name=image)

    prs.save(str(OUT))
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()

