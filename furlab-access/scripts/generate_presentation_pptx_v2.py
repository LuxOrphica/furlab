from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "presentation_assets"
OUT = ROOT / "docs" / "FurLab_DB_Presentation_v2_clean_s.pptx"
IMG_X = Inches(6.45)
IMG_Y = Inches(0.95)
IMG_W = Inches(6.35)
IMG_H = Inches(5.55)


def style_run(run, size=24, bold=False, color=(25, 35, 60)):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Calibri"


def normalize_arrows(text: str) -> str:
    return text.replace("->", "→")


def cleanup_plain_text(text: str) -> str:
    # After extracting code terms into chips, remove dangling arrows before punctuation.
    t = text
    t = re.sub(r"→\s*(?=[,.;:])", "", t)
    t = re.sub(r"(,\s*){2,}", ", ", t)
    t = re.sub(r"\s*,\s*,", ", ", t)
    t = re.sub(r"\(\s*[,;\s]*\)", "", t)
    t = re.sub(r",\s*\)", ")", t)
    t = re.sub(r"\(\s*,\s*", "(", t)
    t = re.sub(r"→\s*$", "", t)
    t = re.sub(r",\s*$", "", t)
    t = re.sub(r":\s*$", "", t)
    t = re.sub(r"\s{2,}", " ", t)
    t = re.sub(r"\s+([,.:;])", r"\1", t)
    return t.strip()


def extract_plain_and_terms(line: str) -> tuple[str, list[str]]:
    line = normalize_arrows(line)
    terms = re.findall(r"`([^`]+)`", line)
    plain = re.sub(r"`[^`]+`", "", line)
    plain = cleanup_plain_text(plain)
    return plain, terms


def add_chip(slide, x: float, y: float, text: str) -> float:
    # Compact "VS Code-like" chip for DB term names.
    w = min(2.45, max(0.85, 0.26 + len(text) * 0.105))
    h = 0.34
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(242, 245, 250)
    shp.line.color.rgb = RGBColor(205, 214, 228)
    shp.line.width = Pt(0.75)
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Consolas"
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(70, 84, 106)
    return w


def estimate_text_box_height(line: str, width_in: float, font_size_pt: int) -> float:
    # Rough but stable estimate to avoid chip overlap on wrapped lines.
    # For Calibri 22-24pt in this deck, ~30 chars per inch is a safe envelope.
    chars_per_inch = 30.0
    max_chars = max(16, int(width_in * chars_per_inch))
    lines = max(1, (len(line) + max_chars - 1) // max_chars)
    line_h = 0.29 if font_size_pt <= 22 else 0.31
    return lines * line_h + 0.08


def add_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12.0), Inches(1.4))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "База данных FurLab"
    style_run(r, size=50, bold=True)

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.0))
    stf = sub.text_frame
    stf.clear()
    p = stf.paragraphs[0]
    r = p.add_run()
    r.text = "Структура, процессы и качество данных • Срез на 11 марта 2026"
    style_run(r, size=24, color=(80, 90, 120))

    slide.shapes.add_shape(1, Inches(0.7), Inches(3.1), Inches(11.8), Inches(0.02))


def add_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    speaker_notes: str,
    image_name: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = normalize_arrows(title)
    style_run(r, size=34, bold=True)

    # Keep a hard gutter before the image column to avoid text overlap.
    left_w = 5.65 if image_name else 12.0
    start_x = 0.6
    y = 1.0
    text_size = 22 if image_name else 24
    for line in bullets:
        plain, terms = extract_plain_and_terms(line)
        if plain:
            plain_h = estimate_text_box_height(plain, left_w, text_size)
            line_box = slide.shapes.add_textbox(Inches(start_x), Inches(y), Inches(left_w), Inches(plain_h))
            ltf = line_box.text_frame
            ltf.clear()
            ltf.word_wrap = True
            p = ltf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = plain
            style_run(run, size=text_size, color=(32, 32, 40))
            y += plain_h + 0.04
        if terms:
            chip_x = start_x + 0.06
            chip_y = y
            for term in terms:
                next_w = min(2.45, max(0.85, 0.26 + len(term) * 0.105))
                if chip_x + next_w > start_x + left_w - 0.12:
                    chip_x = start_x + 0.06
                    chip_y += 0.40
                used_w = add_chip(slide, chip_x, chip_y, term)
                chip_x += used_w + 0.08
            y = chip_y + 0.50
        y += 0.08

    if image_name:
        slide.shapes.add_picture(str(ASSETS / image_name), IMG_X, IMG_Y, width=IMG_W, height=IMG_H)

    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12.0), Inches(0.4))
    ftf = footer.text_frame
    ftf.clear()
    p = ftf.paragraphs[0]
    run = p.add_run()
    run.text = normalize_arrows("Источник: FURLAB / docs / schema + operation contract + Приложение Y")
    style_run(run, size=14, color=(120, 125, 140))

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = normalize_arrows(speaker_notes)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title(prs)

    slides = [
        {
            "title": "1. Зачем нужна эта база",
            "bullets": [
                "Карточка куска хранится в отдельной записи `ScrapPiece`.",
                "Настройки выкладки и запуск расчета фиксируются в `Layout` и `LayoutRun`.",
                "Видим не только текущее состояние, но и историю использования каждого куска.",
            ],
            "notes": "База данных FurLab решает три задачи: учет кусков, поддержка сценариев выкладки и восстановимость решений. Ключевая ценность — мы фиксируем не только текущий статус куска, но и историю применения в конкретном запуске.",
            "image": None,
        },
        {
            "title": "2. На чем основана модель",
            "bullets": [
                "Термины и канон полей взяты из Приложения X.",
                "Сценарии выкладки и параметры — из Приложения Y.",
                "В БД это реализовано через миграции схемы `sql/00x_*.sql`.",
            ],
            "notes": "Модель данных не придумана в коде. Она выведена из приложений X и Y и затем реализована в Access-схеме и миграциях.",
            "image": None,
        },
        {
            "title": "3. Архитектура решения",
            "bullets": [
                "Хранилище данных — файл Access: `.accdb`.",
                "Сервер API принимает запросы от UI и форм через `tools/ui_lab_server.js`.",
                "Формы и UI работают через API, а не напрямую с БД.",
            ],
            "notes": "База отделена от HTTP-слоя. Это позволяет централизованно держать правила валидации и не дублировать логику в интерфейсе.",
            "image": "01_architecture_overview.png",
        },
        {
            "title": "4. Каркас предметной модели",
            "bullets": [
                "В структуре изделия есть деталь, зона и фрагмент: `Part`, `Zone`, `Fragment`.",
                "Схема выкладки хранится отдельно от конкретного запуска: `Layout`, `LayoutRun`.",
                "Складская часть включает карточку куска и место хранения: `ScrapPiece`, `StorageLocation`.",
                "Применение куска к фрагменту фиксируется отдельным фактом `LayoutRunScrapPlacement`.",
            ],
            "notes": "Есть две линии: геометрия изделия и склад. Они сходятся в факте размещения, где видно какой кусок применен, в каком запуске и к какому фрагменту.",
            "image": "02_er_core.png",
        },
        {
            "title": "5. Как устроен учет куска",
            "bullets": [
                "Для каждого куска хранится отдельная карточка `ScrapPiece`.",
                "Качество и текущий статус ведутся как отдельные атрибуты `scrapQuality`, `scrapStatus`.",
                "Геометрия хранится в параметрах `areaMm2`, `maxSpanMm`, `napDirectionDeg`.",
                "Связь с физической биркой обеспечивается инвентарным номером `inventoryTag`.",
            ],
            "notes": "ScrapPiece — центральная таблица учета остатков. Важный инвариант: inventoryTag уникален и однозначно связывает запись с физическим куском.",
            "image": "06_scrappiece_table_sample.png",
        },
        {
            "title": "6. Операции и статусы",
            "bullets": [
                "Статус куска проходит жизненный цикл: доступен, зарезервирован, использован, списан (`Available`, `Reserved`, `Used`, `Discarded`).",
                "Каждое изменение фиксируется операциями: резерв, снятие резерва, использование, списание (`Reserve`, `Release`, `Use`, `Discard`).",
                "Переходы контролируются едиными правилами сервера `status_rules.js`.",
                "История операций сохраняется для аудита в `ScrapTransaction`.",
            ],
            "notes": "Статусы управляются явными правилами. Например reserve разрешен только из Available, а Discarded — терминальное состояние.",
            "image": "04_status_transitions_rules.png",
        },
        {
            "title": "7. Типы выкладок (Приложение Y)",
            "bullets": [
                "Регулярная выкладка по шаблону: `RegularLayout`.",
                "Нерегулярная выкладка по заданным контурам: `IrregularLayout`.",
                "Заполнение остаточной области: `FillRemainingAreaLayout`.",
                "Подбор и размещение складских кусков: `InventoryLayout`.",
            ],
            "notes": "Тип выкладки хранится в Layout.layoutType. Для всех типов используется единая логика хранения параметров и снимка запуска.",
            "image": None,
        },
        {
            "title": "8. Почему результат можно воспроизвести",
            "bullets": [
                "Текущие настройки выкладки сохраняются в `Layout.params`.",
                "Для каждого запуска сохраняется снимок параметров `LayoutRun.paramsSnapshot`.",
                "Результат запуска фиксируется в `LayoutRun.resultSnapshot`.",
                "После применения куска сохраняется итоговый контур `resultContourSnapshot`.",
            ],
            "notes": "Даже если настройки позже изменились, снимок запуска позволяет восстановить фактические входные данные конкретного расчета.",
            "image": "05_traceability_chain.png",
        },
        {
            "title": "9. Два режима инвентарной выкладки",
            "bullets": [
                "Режим A: фрагменты строятся как производные от кусков.",
                "Режим B: куски назначаются на уже готовые фрагменты.",
                "Параметры фильтрации и ограничений задаются конфигурацией `InventoryLayoutConfig`.",
            ],
            "notes": "Оба режима фиксируются в одной и той же таблице размещения, поэтому отчетность остается единой.",
            "image": None,
        },
        {
            "title": "10. Что видно в текущем срезе БД",
            "bullets": [
                "Всего прикладных таблиц: 19.",
                "Карточек кусков (таблица `ScrapPiece`): 121.",
                "По статусам: Available 82, Reserved 28, Used 7, Discarded 4.",
                "Операций в журнале (`ScrapTransaction`): 34.",
                "Фактов размещения (`LayoutRunScrapPlacement`): 7.",
            ],
            "notes": "Складской контур уже рабочий. Одновременно видно зону роста: увеличить полноту фиксации операций Use/Discard.",
            "image": "03_status_distribution.png",
        },
        {
            "title": "11. Эксплуатация и надежность",
            "bullets": [
                "Изменения схемы фиксируются в журнале миграций `SchemaMigrations`.",
                "Для защиты данных есть скрипты `tools/db/*`.",
                "Контракт API проверяется автоматическими тестами `tools/tests/*`.",
                "Операционный регламент вынесен в runbook `docs/RUNBOOK.md`.",
            ],
            "notes": "База сопровождается как продукт: журнал миграций, аварийные процедуры и тесты позволяют безопасно развивать систему.",
            "image": "08_migrations_and_table_counts.png",
        },
        {
            "title": "12. Риски и ближайшие шаги",
            "bullets": [
                "Нужно повысить полноту фиксации жизненного цикла куска.",
                "Нужно насыщать данными настройки подбора в `InventoryLayoutConfig`.",
                "Нужно снижать долю пустого материала в поле `materialId`.",
                "Следующий шаг: KPI качества данных и регулярные отчеты.",
            ],
            "notes": "Переходим от просто работающей БД к управляемой системе с измеримыми показателями качества данных.",
            "image": None,
        },
        {
            "title": "13. Резюме",
            "bullets": [
                "Модель данных соответствует логике Приложения Y.",
                "Есть прозрачная трассируемость и воспроизводимость запусков.",
                "Платформа готова к масштабированию и аналитике качества.",
            ],
            "notes": "БД FurLab уже закрывает ядро задач учета и трассировки. Основной фокус следующего этапа — дисциплина данных и KPI.",
            "image": "07_transaction_history_sample.png",
        },
    ]

    for s in slides:
        add_slide(prs, s["title"], s["bullets"], s["notes"], s["image"])

    prs.save(str(OUT))
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
